"""Polished PowerPoint deck via python-pptx.

Slides are composed from blank layouts rather than the stock title/bullet
placeholders: the default template produces a generic bulleted deck, and the
point of this is a branded one. Charts are native pptx charts so they remain
editable in PowerPoint.
"""

from __future__ import annotations

from io import BytesIO
from typing import Any, Optional

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from services.reports import branding as B
from services.reports.report_data import ReportData

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.6)


def _rgb(hex_string: str) -> RGBColor:
    return RGBColor.from_string(hex_string)


def _blank(prs: Presentation) -> Any:
    return prs.slides.add_slide(prs.slide_layouts[6])


def _band(slide: Any, title: str, subtitle: str = "") -> None:
    """Header band with the slide title."""
    bar = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(1.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(B.INK)
    bar.line.fill.background()
    bar.shadow.inherit = False

    accent = slide.shapes.add_shape(1, 0, Inches(1.0), SLIDE_W, Inches(0.06))
    accent.fill.solid()
    accent.fill.fore_color.rgb = _rgb(B.ACCENT)
    accent.line.fill.background()
    accent.shadow.inherit = False

    box = slide.shapes.add_textbox(MARGIN, Inches(0.16), SLIDE_W - 2 * MARGIN, Inches(0.7))
    frame = box.text_frame
    frame.word_wrap = True
    para = frame.paragraphs[0]
    run = para.add_run()
    run.text = title
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = _rgb("FFFFFF")
    if subtitle:
        sub = frame.add_paragraph()
        srun = sub.add_run()
        srun.text = subtitle
        srun.font.size = Pt(10)
        srun.font.color.rgb = _rgb(B.RULE)


def _footer(slide: Any, page: Optional[int] = None) -> None:
    box = slide.shapes.add_textbox(
        MARGIN, SLIDE_H - Inches(0.45), SLIDE_W - 2 * MARGIN, Inches(0.3)
    )
    para = box.text_frame.paragraphs[0]
    run = para.add_run()
    run.text = B.FOOTER_NOTE
    run.font.size = Pt(7.5)
    run.font.color.rgb = _rgb(B.MUTED)
    if page is not None:
        para.alignment = PP_ALIGN.LEFT


def _title_slide(prs: Presentation, data: ReportData) -> None:
    slide = _blank(prs)
    bg = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = _rgb(B.INK)
    bg.line.fill.background()
    bg.shadow.inherit = False

    accent = slide.shapes.add_shape(1, MARGIN, Inches(2.5), Inches(1.4), Inches(0.09))
    accent.fill.solid()
    accent.fill.fore_color.rgb = _rgb(B.ACCENT)
    accent.line.fill.background()
    accent.shadow.inherit = False

    box = slide.shapes.add_textbox(MARGIN, Inches(2.8), SLIDE_W - 2 * MARGIN, Inches(2.0))
    frame = box.text_frame
    frame.word_wrap = True
    run = frame.paragraphs[0].add_run()
    run.text = data.title
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = _rgb("FFFFFF")

    sub = frame.add_paragraph()
    srun = sub.add_run()
    srun.text = data.subtitle
    srun.font.size = Pt(13)
    srun.font.color.rgb = _rgb(B.RULE)

    tag = slide.shapes.add_textbox(MARGIN, SLIDE_H - Inches(1.0), SLIDE_W - 2 * MARGIN, Inches(0.5))
    trun = tag.text_frame.paragraphs[0].add_run()
    trun.text = "MF Analysis Tool · " + B.FOOTER_NOTE
    trun.font.size = Pt(8)
    trun.font.color.rgb = _rgb(B.MUTED)


def _kpi_slide(prs: Presentation, data: ReportData) -> None:
    slide = _blank(prs)
    _band(slide, "Portfolio at a glance", data.subtitle)

    if data.has_warning:
        panel = slide.shapes.add_shape(1, MARGIN, Inches(1.3), SLIDE_W - 2 * MARGIN, Inches(0.75))
        panel.fill.solid()
        panel.fill.fore_color.rgb = _rgb(B.WARNING_FILL)
        panel.line.color.rgb = _rgb(B.WARNING_INK)
        panel.shadow.inherit = False
        tf = panel.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        run = tf.paragraphs[0].add_run()
        run.text = f"Data quality notice — {data.data_warning}"
        run.font.size = Pt(9)
        run.font.color.rgb = _rgb(B.WARNING_INK)
        top = Inches(2.3)
    else:
        top = Inches(1.7)

    card_w = (SLIDE_W - 2 * MARGIN - Inches(0.4)) / 3
    card_h = Inches(1.5)
    for i, kpi in enumerate(data.kpis[:6]):
        left = MARGIN + (i % 3) * (card_w + Inches(0.2))
        row_top = top + (i // 3) * (card_h + Inches(0.25))
        card = slide.shapes.add_shape(1, left, row_top, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = _rgb(B.BAND)
        card.line.color.rgb = _rgb(B.RULE)
        card.shadow.inherit = False

        frame = card.text_frame
        frame.word_wrap = True
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p0 = frame.paragraphs[0]
        p0.alignment = PP_ALIGN.CENTER
        r0 = p0.add_run()
        r0.text = kpi.label.upper()
        r0.font.size = Pt(9)
        r0.font.color.rgb = _rgb(B.MUTED)

        p1 = frame.add_paragraph()
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = kpi.value
        r1.font.size = Pt(26)
        r1.font.bold = True
        r1.font.color.rgb = _rgb(B.tone_colour(kpi.tone))

        if kpi.caption:
            p2 = frame.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            r2 = p2.add_run()
            r2.text = kpi.caption
            r2.font.size = Pt(9)
            r2.font.color.rgb = _rgb(B.MUTED)
    _footer(slide)


def _chart_slide(prs: Presentation, title: str, mapping: dict[str, float], subtitle: str) -> None:
    slide = _blank(prs)
    _band(slide, title, subtitle)

    chart_data = CategoryChartData()
    chart_data.categories = list(mapping.keys())
    chart_data.add_series("Weight %", tuple(mapping.values()))

    graphic = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT,
        MARGIN,
        Inches(1.5),
        Inches(7.2),
        Inches(5.2),
        chart_data,
    )
    chart = graphic.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(10)

    # Weight table beside the chart, since a doughnut alone is hard to read off.
    rows = min(len(mapping), 10) + 1
    table_shape = slide.shapes.add_table(
        rows, 2, Inches(8.0), Inches(1.6), Inches(4.7), Inches(0.4 * rows)
    )
    table = table_shape.table
    table.columns[0].width = Inches(3.3)
    table.columns[1].width = Inches(1.4)
    for col, label in enumerate(("Bucket", "Weight")):
        cell = table.cell(0, col)
        cell.text = label
        para = cell.text_frame.paragraphs[0]
        para.runs[0].font.size = Pt(11)
        para.runs[0].font.bold = True
    for i, (bucket, weight) in enumerate(list(mapping.items())[:10], start=1):
        table.cell(i, 0).text = str(bucket)[:38]
        table.cell(i, 1).text = f"{weight:.1f}%"
        for col in range(2):
            para = table.cell(i, col).text_frame.paragraphs[0]
            if para.runs:
                para.runs[0].font.size = Pt(10)
    _footer(slide)


def _table_slide(prs: Presentation, title: str, headers: list[str], rows: list[list[str]], subtitle: str) -> None:
    slide = _blank(prs)
    _band(slide, title, subtitle)

    shown = rows[:12]
    shape = slide.shapes.add_table(
        len(shown) + 1, len(headers), MARGIN, Inches(1.5),
        SLIDE_W - 2 * MARGIN, Inches(0.34 * (len(shown) + 1)),
    )
    table = shape.table
    for col, head in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = head
        run = cell.text_frame.paragraphs[0].runs[0]
        run.font.size = Pt(10)
        run.font.bold = True
    for r, record in enumerate(shown, start=1):
        for c, value in enumerate(record):
            cell = table.cell(r, c)
            cell.text = str(value)
            para = cell.text_frame.paragraphs[0]
            if para.runs:
                para.runs[0].font.size = Pt(9)
                if c:
                    para.alignment = PP_ALIGN.RIGHT

    if len(rows) > len(shown):
        note = slide.shapes.add_textbox(
            MARGIN, Inches(1.5) + Inches(0.34 * (len(shown) + 1)) + Inches(0.1),
            SLIDE_W - 2 * MARGIN, Inches(0.3),
        )
        run = note.text_frame.paragraphs[0].add_run()
        run.text = f"Showing the largest {len(shown)} of {len(rows)} — full list in the Excel workbook."
        run.font.size = Pt(9)
        run.font.color.rgb = _rgb(B.MUTED)
    _footer(slide)


def _text_slide(prs: Presentation, title: str, body: str, subtitle: str) -> None:
    slide = _blank(prs)
    _band(slide, title, subtitle)
    box = slide.shapes.add_textbox(MARGIN, Inches(1.6), SLIDE_W - 2 * MARGIN, Inches(5.0))
    frame = box.text_frame
    frame.word_wrap = True
    for i, line in enumerate([ln for ln in body.split("\n") if ln.strip()]):
        para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        run = para.add_run()
        run.text = line.strip()
        run.font.size = Pt(12)
        run.font.color.rgb = _rgb(B.INK)
        para.space_after = Pt(6)
    _footer(slide)


def build_presentation(data: ReportData) -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    _title_slide(prs, data)
    _kpi_slide(prs, data)

    if data.risk.rows:
        _table_slide(
            prs, "Risk & return", ["Metric", "Value"],
            [[k, v] for k, v in data.risk.rows], data.subtitle,
        )

    for name, mapping in data.allocations.items():
        _chart_slide(prs, name, mapping, data.subtitle)

    if data.top_holdings:
        _table_slide(
            prs, "Largest look-through holdings", ["Security", "Weight"],
            [[n, f"{w:.2f}%"] for n, w in data.top_holdings], data.subtitle,
        )

    if data.overlap_rows:
        _table_slide(
            prs, "Fund overlap", ["Measure", "Value"],
            [[k, v] for k, v in data.overlap_rows], data.subtitle,
        )

    if data.holdings_display:
        headers = ["Scheme", "Current value", "Return", "Weight"]
        rows = [
            [r.get("Scheme", ""), r.get("Current value", ""), r.get("Return", ""), r.get("Weight", "")]
            for r in data.holdings_display
        ]
        _table_slide(prs, "Holdings", headers, rows, data.subtitle)

    if data.ai_summary:
        _text_slide(prs, "AI review", data.ai_summary, data.subtitle)

    return prs


def pptx_bytes(data: ReportData) -> bytes:
    buffer = BytesIO()
    build_presentation(data).save(buffer)
    return buffer.getvalue()
