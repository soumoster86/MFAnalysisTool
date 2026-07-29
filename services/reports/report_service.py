"""PDF / Excel / PowerPoint report generation."""

from __future__ import annotations

from datetime import datetime
from io import BytesIO
from pathlib import Path
from typing import Any, Optional

import pandas as pd

from config.settings import PROJECT_ROOT, settings
from utils.logging_config import get_logger

logger = get_logger(__name__)


class ReportService:
    """Generate downloadable portfolio / fund reports."""

    def __init__(self, output_dir: Optional[Path] = None) -> None:
        self.output_dir = Path(output_dir or (PROJECT_ROOT / "data" / "reports"))
        self.output_dir.mkdir(parents=True, exist_ok=True)

    def generate_excel(
        self,
        *,
        portfolio_rows: list[dict[str, Any]],
        metrics: Optional[dict[str, Any]] = None,
        health: Optional[dict[str, Any]] = None,
        filename: Optional[str] = None,
    ) -> Path:
        fname = filename or f"mf_report_{datetime.now():%Y%m%d_%H%M%S}.xlsx"
        path = self.output_dir / fname
        with pd.ExcelWriter(path, engine="openpyxl") as writer:
            pd.DataFrame(portfolio_rows).to_excel(writer, sheet_name="Portfolio", index=False)
            if metrics:
                pd.DataFrame([metrics]).to_excel(writer, sheet_name="Metrics", index=False)
            if health:
                pd.DataFrame([health]).to_excel(writer, sheet_name="Health", index=False)
        logger.info("Excel report written to {}", path)
        return path

    def generate_pdf(
        self,
        *,
        title: str,
        summary_lines: list[str],
        table_rows: Optional[list[dict[str, Any]]] = None,
        ai_summary: str = "",
        filename: Optional[str] = None,
    ) -> Path:
        from reportlab.lib import colors
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

        fname = filename or f"mf_report_{datetime.now():%Y%m%d_%H%M%S}.pdf"
        path = self.output_dir / fname
        doc = SimpleDocTemplate(str(path), pagesize=A4)
        styles = getSampleStyleSheet()
        story = [
            Paragraph(title, styles["Title"]),
            Paragraph(f"Generated: {datetime.now():%Y-%m-%d %H:%M}", styles["Normal"]),
            Spacer(1, 12),
        ]
        for line in summary_lines:
            story.append(Paragraph(line, styles["BodyText"]))
            story.append(Spacer(1, 6))
        if table_rows:
            headers = list(table_rows[0].keys())
            data = [headers] + [[str(r.get(h, "")) for h in headers] for r in table_rows]
            t = Table(data, repeatRows=1)
            t.setStyle(
                TableStyle(
                    [
                        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1a1a2e")),
                        ("TEXTCOLOR", (0, 0), (-1, 0), colors.whitesmoke),
                        ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
                        ("FONTSIZE", (0, 0), (-1, -1), 8),
                        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.HexColor("#f5f5f5"), colors.white]),
                    ]
                )
            )
            story.append(Spacer(1, 12))
            story.append(t)
        if ai_summary:
            story.append(Spacer(1, 16))
            story.append(Paragraph("AI Summary", styles["Heading2"]))
            story.append(Paragraph(ai_summary.replace("\n", "<br/>"), styles["BodyText"]))
        story.append(Spacer(1, 20))
        story.append(
            Paragraph(
                "Disclaimer: Educational report only. Not investment advice.",
                styles["Italic"],
            )
        )
        doc.build(story)
        logger.info("PDF report written to {}", path)
        return path

    def generate_pptx(
        self,
        *,
        title: str,
        bullets: list[str],
        metrics: Optional[dict[str, Any]] = None,
        filename: Optional[str] = None,
    ) -> Path:
        from pptx import Presentation
        from pptx.util import Inches, Pt

        fname = filename or f"mf_report_{datetime.now():%Y%m%d_%H%M%S}.pptx"
        path = self.output_dir / fname
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        slide.placeholders[1].text = f"MF Analysis Tool · {datetime.now():%Y-%m-%d}"

        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        slide2.shapes.title.text = "Key Points"
        body = slide2.shapes.placeholders[1].text_frame
        body.clear()
        for i, b in enumerate(bullets):
            if i == 0:
                body.text = b
            else:
                p = body.add_paragraph()
                p.text = b
                p.level = 0

        if metrics:
            slide3 = prs.slides.add_slide(prs.slide_layouts[1])
            slide3.shapes.title.text = "Metrics Snapshot"
            tf = slide3.shapes.placeholders[1].text_frame
            tf.clear()
            for i, (k, v) in enumerate(metrics.items()):
                line = f"{k}: {v}"
                if i == 0:
                    tf.text = line
                else:
                    p = tf.add_paragraph()
                    p.text = line

        prs.save(str(path))
        logger.info("PPTX report written to {}", path)
        return path

    def excel_bytes(self, portfolio_rows: list[dict[str, Any]]) -> bytes:
        buf = BytesIO()
        with pd.ExcelWriter(buf, engine="openpyxl") as writer:
            pd.DataFrame(portfolio_rows).to_excel(writer, index=False)
        return buf.getvalue()
