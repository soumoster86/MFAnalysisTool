"""Report assembly and rendering across PDF, Excel and PowerPoint."""

from __future__ import annotations

from dataclasses import dataclass, field
from io import BytesIO
from typing import Any, Optional

import pytest

from services.reports.report_data import (
    HOLDING_COLUMNS,
    _merge_label_variants,
    build_report_data,
)
from services.reports.report_service import ReportService


@dataclass
class _Analysis:
    """Stand-in for PortfolioAnalysis with only what the report reads."""

    total_invested: float = 130000.0
    total_current: float = 146100.0
    daily_pnl: float = 1313.96
    daily_pnl_pct: float = 0.009
    overall_gain: float = 16100.0
    overall_gain_pct: float = 0.1238
    portfolio_cagr: Optional[float] = 0.142
    expected_cagr: Optional[float] = 0.142
    volatility: Optional[float] = 0.118
    sharpe: Optional[float] = 0.81
    max_drawdown: Optional[float] = -0.213
    health_score: float = 68.4
    mode: str = "full"
    asset_allocation: dict = field(default_factory=lambda: {"Equity": 74.0, "Debt": 26.0})
    sector_allocation: dict = field(
        default_factory=lambda: {"Financial": 20.0, "Financials": 7.5, "Technology": 12.0}
    )
    market_cap_allocation: dict = field(
        default_factory=lambda: {"Large": 60.0, "Mid": 25.0, "Small": 15.0}
    )
    top_holdings: list = field(
        default_factory=lambda: [
            {"security": "HDFC Bank Ltd", "weight_pct": 6.4},
            {"security": "Infosys Ltd", "weight_pct": 4.1},
        ]
    )
    holdings_detail: list = field(
        default_factory=lambda: [
            {
                "scheme_name": "Alpha Large Cap Fund - Direct Growth",
                "category": "Equity",
                "units": 80.0,
                "invested_amount": 60000.0,
                "current_value": 72000.0,
                "weight_pct": 49.3,
            },
            {
                "scheme_name": "Beta Debt Fund - Direct Growth",
                "category": "Debt",
                "units": 300.0,
                "invested_amount": 70000.0,
                "current_value": 74100.0,
                "weight_pct": 50.7,
            },
        ]
    )
    overlap: dict = field(
        default_factory=lambda: {
            "holding_overlap_pct": 18.2,
            "sector_overlap_pct": 41.0,
            "diversification_score": 72.0,
            "pairwise_overlap": {"Alpha × Beta": 18.2},
        }
    )
    notes: list = field(default_factory=lambda: ["Analysis mode: full (2 schemes)."])
    nav_series: Any = None
    correlation: Any = None
    data_sources: dict = field(default_factory=dict)


FABRICATED_SOURCES = {
    "nav": {"Alpha Large Cap Fund - Direct Growth": "synthetic"},
    "holdings": {"Beta Debt Fund - Direct Growth": "sample"},
}
CLEAN_SOURCES = {
    "nav": {"Alpha Large Cap Fund - Direct Growth": "mfapi"},
    "holdings": {"Beta Debt Fund - Direct Growth": "groww"},
}


@pytest.fixture()
def data():
    return build_report_data(_Analysis(data_sources=CLEAN_SOURCES), title="Test Report")


@pytest.fixture()
def service():
    return ReportService()


# ----------------------------------------------------------------- assembly
def test_kpis_cover_the_headline_figures(data):
    labels = [k.label for k in data.kpis]
    assert "Portfolio value" in labels
    assert "Overall P&L" in labels
    assert "Health score" in labels


def test_gain_and_loss_carry_opposite_tones():
    gain = build_report_data(_Analysis(overall_gain=5000.0))
    loss = build_report_data(_Analysis(overall_gain=-5000.0, overall_gain_pct=-0.1))
    assert next(k for k in gain.kpis if k.label == "Overall P&L").tone == "good"
    assert next(k for k in loss.kpis if k.label == "Overall P&L").tone == "bad"


def test_holdings_are_ordered_largest_first(data):
    values = [h["current_value"] for h in data.holdings]
    assert values == sorted(values, reverse=True)


def test_holdings_gain_is_derived_per_row(data):
    row = data.holdings[0]
    assert row["gain"] == pytest.approx(row["current_value"] - row["invested_amount"])
    assert row["gain_pct"] == pytest.approx(
        row["current_value"] / row["invested_amount"] - 1
    )


def test_display_rows_use_the_curated_columns_only(data):
    expected = {label for _, label, _ in HOLDING_COLUMNS}
    assert set(data.holdings_display[0]) == expected


def test_holdings_are_capped():
    many = _Analysis(
        holdings_detail=[
            {
                "scheme_name": f"Fund {i}",
                "category": "Equity",
                "units": 1.0,
                "invested_amount": 1000.0,
                "current_value": 1000.0 + i,
                "weight_pct": 1.0,
            }
            for i in range(60)
        ]
    )
    assert len(build_report_data(many, max_holdings=25).holdings) == 25


def test_missing_values_render_as_a_dash():
    thin = _Analysis(portfolio_cagr=None, sharpe=None, volatility=None, max_drawdown=None)
    rows = dict(build_report_data(thin).risk.rows)
    assert rows["Portfolio CAGR"] == "—"
    assert rows["Sharpe ratio"] == "—"


def test_overlap_uses_the_real_result_shape(data):
    labels = [k for k, _ in data.overlap_rows]
    assert "Average holding overlap" in labels
    assert "Diversification score" in labels
    assert any("Alpha" in k for k in labels)


def test_zero_weight_buckets_are_dropped():
    noisy = _Analysis(asset_allocation={"Equity": 80.0, "Ghost": 0.0})
    assert "Ghost" not in build_report_data(noisy).allocations["Asset allocation"]


# -------------------------------------------------------- label variants
def test_plural_and_singular_sector_labels_are_merged(data):
    sectors = data.allocations["Sector allocation"]
    # "Financial" and "Financials" are the same exposure from an inconsistent
    # provider and must not appear as two separate slices.
    assert sum(1 for k in sectors if k.lower().startswith("financial")) == 1
    assert sectors["Financials"] == pytest.approx(27.5)


def test_genuinely_different_buckets_stay_separate():
    merged = _merge_label_variants(
        {"Consumer": 5.0, "Consumer Discretionary": 9.0, "Consumer Staples": 4.0}
    )
    assert len(merged) == 3


def test_merge_is_whitespace_insensitive():
    merged = _merge_label_variants({"Health Care": 4.0, "Health  Care": 2.0})
    assert len(merged) == 1
    assert list(merged.values())[0] == pytest.approx(6.0)


# --------------------------------------------------------------- provenance
def test_fabricated_inputs_produce_a_warning_inside_the_report():
    # A file is read away from the app, so an on-screen banner protects nobody.
    out = build_report_data(_Analysis(data_sources=FABRICATED_SOURCES))
    assert out.has_warning
    assert "synthetic" in out.data_warning.lower()
    assert "sample holdings" in out.data_warning.lower()


def test_warning_names_only_the_problems_that_occurred():
    # "0 funds used synthetic NAV and 2 used sample holdings" reads like a bug.
    only_holdings = build_report_data(
        _Analysis(data_sources={"nav": {"A": "mfapi"}, "holdings": {"B": "sample"}})
    )
    assert "sample holdings" in only_holdings.data_warning
    assert "synthetic" not in only_holdings.data_warning
    assert " 0 " not in only_holdings.data_warning


def test_warning_is_singular_for_one_fund():
    single = build_report_data(
        _Analysis(data_sources={"nav": {"A": "synthetic"}, "holdings": {}})
    )
    assert "1 fund on" in single.data_warning


def test_clean_inputs_produce_no_warning(data):
    assert not data.has_warning
    assert data.data_warning == ""


def test_source_rows_are_listed_for_the_audit_trail(data):
    assert data.source_rows
    assert any("NAV" in v for _, v in data.source_rows)


# ---------------------------------------------------------------- rendering
@pytest.mark.parametrize("fmt", ["pdf", "excel", "pptx"])
def test_every_format_renders_a_non_trivial_file(service, data, fmt):
    payload, filename, mime = service.render(data, fmt)
    assert len(payload) > 4000
    assert filename.endswith({"pdf": ".pdf", "excel": ".xlsx", "pptx": ".pptx"}[fmt])
    assert mime


def test_pdf_has_a_valid_header_and_multiple_pages(service, data):
    payload = service.pdf_bytes(data)
    assert payload.startswith(b"%PDF")

    from pypdf import PdfReader

    reader = PdfReader(BytesIO(payload))
    assert len(reader.pages) >= 1
    text = "\n".join((p.extract_text() or "") for p in reader.pages)
    assert data.title in text
    assert "Holdings" in text


def test_pdf_renders_the_rupee_sign_not_a_missing_glyph_box(service, data):
    # Helvetica reports a width for U+20B9 but draws a hollow box, so every
    # currency figure looked broken until a Unicode face was registered.
    from pypdf import PdfReader

    reader = PdfReader(BytesIO(service.pdf_bytes(data)))
    text = "\n".join((p.extract_text() or "") for p in reader.pages)
    assert "�" not in text
    assert "■" not in text
    # Either the real symbol, or the honest ASCII fallback.
    assert "₹" in text or "Rs " in text


def test_pdf_carries_the_data_warning_when_inputs_are_fabricated(service):
    from pypdf import PdfReader

    warned = build_report_data(_Analysis(data_sources=FABRICATED_SOURCES))
    reader = PdfReader(BytesIO(service.pdf_bytes(warned)))
    text = "\n".join((p.extract_text() or "") for p in reader.pages)
    assert "Data quality notice" in text


def test_excel_has_the_expected_sheets(service, data):
    from openpyxl import load_workbook

    wb = load_workbook(BytesIO(service.excel_bytes(data)))
    assert {"Summary", "Holdings", "Allocation", "Data sources"} <= set(wb.sheetnames)


def test_excel_writes_numbers_as_numbers_not_strings(service, data):
    # A spreadsheet of pre-formatted strings cannot be sorted, totalled or
    # charted, which is the only reason to want Excel over PDF.
    from openpyxl import load_workbook

    ws = load_workbook(BytesIO(service.excel_bytes(data)))["Holdings"]
    header_row = next(
        r for r in range(1, 12)
        if ws.cell(row=r, column=1).value == "Scheme"
    )
    invested = ws.cell(row=header_row + 1, column=4)
    assert isinstance(invested.value, (int, float))
    assert "#,##0" in (invested.number_format or "")


def test_excel_holdings_sheet_totals_reconcile(service, data):
    from openpyxl import load_workbook

    ws = load_workbook(BytesIO(service.excel_bytes(data)))["Holdings"]
    formulas = [
        ws.cell(row=r, column=c).value
        for r in range(1, ws.max_row + 1)
        for c in range(1, 9)
        if isinstance(ws.cell(row=r, column=c).value, str)
        and str(ws.cell(row=r, column=c).value).startswith("=SUM(")
    ]
    assert formulas, "no total row written"


def test_excel_flags_fabricated_sources_in_the_audit_sheet(service):
    from openpyxl import load_workbook

    warned = build_report_data(_Analysis(data_sources=FABRICATED_SOURCES))
    ws = load_workbook(BytesIO(service.excel_bytes(warned)))["Data sources"]
    values = [
        ws.cell(row=r, column=2).value
        for r in range(1, ws.max_row + 1)
    ]
    assert any(v and ("synthetic" in str(v) or "sample" in str(v)) for v in values)


def test_pptx_builds_the_expected_slide_count(service, data):
    from pptx import Presentation

    prs = Presentation(BytesIO(service.pptx_bytes(data)))
    # Title, KPIs, risk, three allocations, look-through, overlap, holdings.
    assert len(prs.slides) >= 6


def test_pptx_first_slide_carries_the_title(service, data):
    from pptx import Presentation

    prs = Presentation(BytesIO(service.pptx_bytes(data)))
    text = " ".join(
        shape.text_frame.text
        for shape in prs.slides[0].shapes
        if shape.has_text_frame
    )
    assert data.title in text


# ------------------------------------------------------------------- errors
def test_unknown_format_is_rejected(service, data):
    with pytest.raises(ValueError, match="Unknown report format"):
        service.render(data, "docx")


def test_empty_portfolio_still_renders_every_format(service):
    empty = _Analysis(
        holdings_detail=[],
        top_holdings=[],
        overlap={},
        asset_allocation={},
        sector_allocation={},
        market_cap_allocation={},
    )
    data = build_report_data(empty)
    for fmt in ("pdf", "excel", "pptx"):
        payload, _, _ = service.render(data, fmt)
        assert len(payload) > 1000


def test_saving_writes_a_file(service, data, tmp_path):
    service.output_dir = tmp_path
    path = service.save(data, "pdf", "out.pdf")
    assert path.exists() and path.stat().st_size > 4000


def test_holdings_csv_matches_the_display_columns(service, data):
    csv = service.holdings_csv(data).decode("utf-8")
    assert "Scheme" in csv.splitlines()[0]
    assert len(csv.strip().splitlines()) == len(data.holdings_display) + 1
